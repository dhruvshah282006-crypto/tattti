from fastapi import FastAPI, UploadFile, File
from fastapi.responses import JSONResponse
from pptx import Presentation
import shutil
import os
import uvicorn

app = FastAPI()

def analyze_pptx(path):
    prs = Presentation(path)
    stats = {
        "slide_count": len(prs.slides),
        "word_count": 0,
        "image_count": 0,
        "chart_count": 0
    }
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                stats["word_count"] += len(shape.text.split())
            if shape.shape_type == 13: # Picture
                stats["image_count"] += 1
            if shape.shape_type == 3:  # Chart
                stats["chart_count"] += 1
    return stats

def calculate_score_out_of_10(stats):
    """Normalized 1-10 Logic"""
    TARGET_SLIDES, TARGET_IMAGES, TARGET_CHARTS = 10, 5, 2
    
    if stats['slide_count'] < 5:
        return 0.0, "TOO_SHORT"

    # Weighted scoring
    slide_score = min(4.0, (stats['slide_count'] / TARGET_SLIDES) * 4.0)
    image_score = min(3.0, (stats['image_count'] / TARGET_IMAGES) * 3.0)
    chart_score = min(3.0, (stats['chart_count'] / TARGET_CHARTS) * 3.0)

    total_score = slide_score + image_score + chart_score
    words_per_slide = stats['word_count'] / stats['slide_count']
    
    if words_per_slide > 60: total_score -= 2.0
    elif words_per_slide < 10: total_score -= 1.0

    return round(max(0.0, min(10.0, total_score)), 1), "READY"

@app.get("/")
def hello():
    return JSONResponse({"message": "backend is running!"})

@app.post("/score-ppt")
async def score_ppt(file: UploadFile = File(...)):
    try:
        temp_path = f"temp_{file.filename}"
        with open(temp_path, "wb") as buffer:
            shutil.copyfileobj(file.file, buffer)

        stats = analyze_pptx(temp_path)
        score, flag = calculate_score_out_of_10(stats)

        os.remove(temp_path)

        return JSONResponse({
            "score": score,
            "status": flag,
            "stats": stats
        })
    except Exception as e:
        return JSONResponse(status_code=500, content={"error": str(e)})

if __name__ == "__main__":
    port = int(os.environ.get("PORT", 8000))
    # Note: changed "main:app" to app so it runs directly if called as __main__
    uvicorn.run(app, host="0.0.0.0", port=port)